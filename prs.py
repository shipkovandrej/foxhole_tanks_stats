import sys
from pptx import Presentation, exc
from pptx.util import Pt

from config import links

from main import get_vehicle_stats

try:
    prs = Presentation('foxhole_tanks_stats.pptx')
except exc.PackageNotFoundError:
    print("\nНет презентации foxhole_tanks_stats.pptx в каталоге\n")
    input("Enter для выхода")
    sys.exit(0)

# Проверка на индексы текста
# slide = prs.slides[3]
# block = slide.shapes[1].text_frame.text
# print(block)
# quit()

for slide in prs.slides:
    text = slide.shapes.title.text
    if text in links:
        tank_name = text
        result = get_vehicle_stats(links[tank_name])

        block = slide.shapes[1].text_frame

        block.text = f"ХП - {result['hp']} ед"

        p = block.add_paragraph()
        p.text = f"Броня - {result['armor']} ед"

        p = block.add_paragraph()
        p.text = f"Подбит при - <{result['disable']}% ХП"

        p = block.add_paragraph()
        p.text = "П.П/У: |"
        for key, value in result['health_table'].items():
            p.text += f" {key} - {value} |"

        p = block.add_paragraph()
        p.text = f"Базовый шанс пробития - {result['min_pin_chance']}%"

        p = block.add_paragraph()
        p.text = f"Максимальная степень износа танковой брони - {result['max_pin_chance']}%"

        p = block.add_paragraph()
        p.text = "Шанс подбития подсистем:"

        for key, value in result['subsystems_pin_chance'].items():
            p = block.add_paragraph()
            p.text = f"{key} - {value}%"
            p.level = 1

        for i in block.paragraphs:
            i.runs[0].font.size = Pt(16)
            i.runs[0].font.name = "Bahnschrift Condensed"

        block = slide.shapes[2].text_frame

        block.text = f"Стоимость полной починки - {result['repair_cost']} бматов"

        p = block.add_paragraph()
        p.text = (f"Скорость: по дороге - {result['speed']['road']} м/с, "
                  f"по внедорожью - {result['speed']['off-road']} м/с")

        p = block.add_paragraph()
        p.text = f"Объём бака: {result['fuel_tank']} л"

        p = block.add_paragraph()
        p.text = f"Потребление: {result['fuel_consumption_rate']} л/мин"

        p = block.add_paragraph()
        p.text = f"Время хода: {result['fuel_autonomy']}"

        p = block.add_paragraph()
        p.text = "Вооружение:"

        for i in result['guns']:
            weapon_name = i
            reload_dict = result['reload'][weapon_name]

            if 'firing_duration' in reload_dict:
                weapon_reload = f"{reload_dict['firing_duration']} + {reload_dict['reload_duration']}"
            else:
                weapon_reload = f"{reload_dict['reload_duration']}"

            weapon_range = result['range'][weapon_name]['range']

            p = block.add_paragraph()
            text = f"{weapon_name} | {weapon_reload} сек. | {weapon_range} метров"
            p.text = text
            p.level = 1

        p = block.add_paragraph()
        p.text = "Экипаж:"

        for i in result['crew']:
            p = block.add_paragraph()
            p.text = i
            p.level = 1

        for i in block.paragraphs:
            i.runs[0].font.size = Pt(16)
            i.runs[0].font.name = "Bahnschrift Condensed"
        prs.save("foxhole_tanks_stats.pptx")

input("Enter для выхода")
